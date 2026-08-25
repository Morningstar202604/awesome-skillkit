#!/usr/bin/env python3
"""Render slides_spec.json into a .pptx deck (ppt-builder skill).

Usage:
    python make_pptx.py <spec.json> <output.pptx>

Spec schema:
    {"deck_title": str,
     "slides": [{"title": str, "bullets": [str, ...], "notes": str}]}

Exit codes:
    0 = wrote output
    2 = spec invalid (details on stderr)
    3 = python-pptx not installed
"""

import json
import sys

MAX_BULLETS = 5


def validate_spec(spec):
    """Return a list of human-readable errors (empty list = valid)."""
    errors = []
    if not isinstance(spec, dict):
        return ["spec root must be an object"]
    if not isinstance(spec.get("deck_title"), str) or not spec["deck_title"].strip():
        errors.append("deck_title must be a non-empty string")
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        errors.append("slides must be a non-empty array")
        return errors
    for i, s in enumerate(slides, 1):
        where = f"slide {i}"
        if not isinstance(s, dict):
            errors.append(f"{where}: must be an object")
            continue
        if not isinstance(s.get("title"), str) or not s["title"].strip():
            errors.append(f"{where}: title must be a non-empty string")
        bullets = s.get("bullets")
        if not isinstance(bullets, list) or not bullets:
            errors.append(f"{where}: bullets must be a non-empty array")
        else:
            if len(bullets) > MAX_BULLETS:
                errors.append(
                    f"{where}: {len(bullets)} bullets exceeds max {MAX_BULLETS}"
                )
            for b in bullets:
                if not isinstance(b, str) or not b.strip():
                    errors.append(f"{where}: every bullet must be a non-empty string")
        if not isinstance(s.get("notes", ""), str):
            errors.append(f"{where}: notes must be a string when present")
    return errors


def render(spec, out_path):
    try:
        from pptx import Presentation
    except ImportError:
        print(
            "ERROR: python-pptx is not installed. Run: pip install python-pptx",
            file=sys.stderr,
        )
        return 3

    prs = Presentation()
    layout_title = prs.slide_layouts[0]
    layout_bullets = prs.slide_layouts[1]

    cover = prs.slides.add_slide(layout_title)
    cover.shapes.title.text = spec["deck_title"]

    for s in spec["slides"]:
        slide = prs.slides.add_slide(layout_bullets)
        slide.shapes.title.text = s["title"]
        body = slide.placeholders[1].text_frame
        for j, bullet in enumerate(s["bullets"]):
            para = body.paragraphs[0] if j == 0 else body.add_paragraph()
            para.text = bullet
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]

    prs.save(out_path)
    return 0


def main(argv):
    if len(argv) != 3:
        print(__doc__, file=sys.stderr)
        return 2
    spec_path, out_path = argv[1], argv[2]
    try:
        with open(spec_path, encoding="utf-8-sig") as f:
            spec = json.load(f)
    except (OSError, ValueError) as e:
        print(f"ERROR: cannot read spec: {e}", file=sys.stderr)
        return 2

    errors = validate_spec(spec)
    if errors:
        for e in errors:
            print(f"SPEC ERROR: {e}", file=sys.stderr)
        return 2

    rc = render(spec, out_path)
    if rc == 0:
        print(f"wrote {out_path} ({len(spec['slides'])} slides + cover)")
    return rc


if __name__ == "__main__":
    sys.exit(main(sys.argv))
